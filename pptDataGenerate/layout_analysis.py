import os
import json
import subprocess
import shutil
from pptx import Presentation
from PIL import Image
from pdf2image import convert_from_path

def emu_to_px(emu):
    if emu is None: return 0
    # 这里的 9144 是一个约数，后续会根据 PDF 实际像素进行 scale 缩放
    return int(emu / 9144)

def ppt_to_pdf_libreoffice(ppt_path, output_dir):
    """单次调用 LibreOffice，将整个 PPT 转换为一个 PDF 文件"""
    print(f"  [Step 1] 正在通过 LibreOffice 转换 PDF...")
    abs_ppt_path = os.path.abspath(ppt_path)
    cmd = [
        'libreoffice', '--headless', '--convert-to', 'pdf',
        '--outdir', output_dir, abs_ppt_path
    ]
    subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    pdf_name = os.path.splitext(os.path.basename(ppt_path))[0] + ".pdf"
    return os.path.join(output_dir, pdf_name)

def crop_and_save(big_img, box, save_path):
    l, t, r, b = box
    l, t = max(0, l), max(0, t)
    r, b = min(big_img.width, r), min(big_img.height, b)
    if r > l and b > t:
        cropped = big_img.crop((l, t, r, b))
        cropped.save(save_path, "PNG")

def get_table_html(shape):
    """将 PPT 表格对象转换为 HTML 字符串"""
    table = shape.table
    html = "<table>"
    for row in table.rows:
        html += "<tr>"
        for cell in row.cells:
            # 简单的单元格合并逻辑处理 (如果有)
            # 这里默认处理基础单元格内容
            html += f"<td>{cell.text_frame.text.strip()}</td>"
        html += "</tr>"
    html += "</table>"
    return html

def process_batch_ppts(input_folder, base_output_path):
    """
    批量处理文件夹下的所有 PPT/PPTX
    """
    os.makedirs(base_output_path, exist_ok=True)
    ppt_files = [f for f in os.listdir(input_folder) if f.lower().endswith(('.pptx', '.ppt'))]
    print(f"找到 {len(ppt_files)} 个待处理文件。")
    index=0
    for ppt_file in ppt_files:
        ppt_path = os.path.join(input_folder, ppt_file)
        ppt_stem = os.path.splitext(ppt_file)[0]
        
        ppt_root_dir = os.path.join(base_output_path, ppt_stem)
        os.makedirs(ppt_root_dir, exist_ok=True)
        
        print(f"\n>>> 正在处理文件: {ppt_file}")
        temp_worker_dir = os.path.join(ppt_root_dir, "temp_worker")
        os.makedirs(temp_worker_dir, exist_ok=True)

        try:
            if ppt_file.lower().endswith('.ppt'):
                print(f"  [Wait] 检测到旧版 PPT，正在预转换为 PPTX...")
                convert_cmd = [
                    'libreoffice', '--headless', '--convert-to', 'pptx',
                    '--outdir', temp_worker_dir, os.path.abspath(ppt_path)
                ]
                subprocess.run(convert_cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                working_ppt_path = os.path.join(temp_worker_dir, ppt_stem + ".pptx")
            else:
                working_ppt_path = ppt_path

            pdf_path = ppt_to_pdf_libreoffice(ppt_path, temp_worker_dir)
            prs = Presentation(working_ppt_path)
            total_slides = len(prs.slides)
            print(f"  [Step 2] 检测到 {total_slides} 页内容，开始逐页解析...")

            for page_index in range(total_slides):
                page_dir = os.path.join(ppt_root_dir, f"page_{page_index}")
                images_dir = os.path.join(page_dir, "images")
                tables_dir = os.path.join(page_dir, "tables")
                os.makedirs(images_dir, exist_ok=True)
                os.makedirs(tables_dir, exist_ok=True)

                images = convert_from_path(
                    pdf_path, dpi=300, 
                    first_page=page_index + 1, 
                    last_page=page_index + 1
                )
                
                if not images: continue
                big_img = images[0]
                big_w, big_h = big_img.size
                big_img.save(os.path.join(page_dir, "full_slide.png"), "PNG")

                scale_x = big_w / emu_to_px(prs.slide_width)
                scale_y = big_h / emu_to_px(prs.slide_height)

                # 最终输出的列表格式
                final_json_elements = []
                img_idx, table_idx = 1, 1
                slide = prs.slides[page_index]
                slide_width_emu = prs.slide_width
                slide_height_emu = prs.slide_height
                for shape in slide.shapes:
                    try:
                        # 获取基本坐标
                        x0 = int((shape.left / slide_width_emu) * big_w)
                        y0 = int((shape.top / slide_height_emu) * big_h)
                        x1 = int(((shape.left + shape.width) / slide_width_emu) * big_w)
                        y1 = int(((shape.top + shape.height) / slide_height_emu) * big_h)
                        
                        bbox = [x0, y0, x1, y1]
                        #x0 = int(emu_to_px(shape.left) * scale_x)
                        #y0 = int(emu_to_px(shape.top) * scale_y)
                        #x1 = x0 + int(emu_to_px(shape.width) * scale_x)
                        #y1 = y0 + int(emu_to_px(shape.height) * scale_y)
                        #bbox = [x0, y0, x1, y1]

                        # --- 1. 表格处理 ---
                        if shape.has_table:
                            save_name = f"slide_{page_index}_table_{table_idx}.jpg"
                            rel_path = f"slide_{page_index}_tables/{save_name}"
                            crop_and_save(big_img, (x0, y0, x1, y1), os.path.join(tables_dir, save_name))
                            
                            final_json_elements.append({
                                "type": "table",
                                "img_path": f"../posterData/mineru_output/{ppt_stem}/visuals/{rel_path}",
                                "table_caption": [],
                                "table_footnote": [],
                                "table_body": get_table_html(shape),
                                "bbox": bbox
                            })
                            table_idx += 1
                            continue 

                        # --- 2. 文本处理 ---
                        if shape.has_text_frame and shape.text.strip():
                            paragraphs = [p for p in shape.text_frame.paragraphs if p.text.strip()]
                            
                            # 默认不认为是列表，除非满足条件
                            is_list = any(p.level > 0 for p in paragraphs) or len(paragraphs) > 1
                            t_level = None

                            # 安全地检查占位符类型
                            if shape.is_placeholder:
                                try:
                                    ph_type = shape.placeholder_format.type
                                    if ph_type in [1, 3]: # 1: Title, 3: Center Title
                                        is_list = False
                                        t_level = 1
                                except (ValueError, AttributeError):
                                    pass # 如果获取占位符详情失败，按普通文本处理

                            if not t_level and paragraphs and paragraphs[0].font.bold:
                                t_level = 1

                            if is_list:
                                final_json_elements.append({
                                    "type": "list",
                                    "sub_type": "text",
                                    "list_items": [p.text.strip() for p in paragraphs],
                                    "bbox": bbox
                                })
                            else:
                                text_item = {
                                    "type": "text",
                                    "text": shape.text.strip(),
                                    "bbox": bbox
                                }
                                if t_level:
                                    text_item["text_level"] = t_level
                                final_json_elements.append(text_item)
                            continue

                        # --- 3. 视觉元素处理 ---
                        shape_type_code = int(shape.shape_type)
                        # 6: Group, 13: Picture, 24: Canvas/SmartArt
                        if shape_type_code in [6, 13, 24] or not shape.has_text_frame:
                            if (x1 - x0) < 5 or (y1 - y0) < 5:
                                continue
                                
                            save_name = f"slide_{page_index}_fig_{img_idx}.jpg"
                            rel_path = f"images/slide_{page_index}_{save_name}"
                            crop_and_save(big_img, (x0, y0, x1, y1), os.path.join(images_dir, save_name))
                            
                            final_json_elements.append({
                                "type": "image",
                                "img_path": f"../posterData/mineru_output/{ppt_stem}/visuals/{rel_path}",
                                "image_caption": [],
                                "image_footnote": [],
                                "bbox": bbox
                            })
                            img_idx += 1
                            
                    except Exception as shape_err:
                        # 打印单个 shape 的错误，但不中断整页解析
                        print(f"      [Warning] 解析第 {page_index} 页的某个元素时跳过: {shape_err}")
                        continue
                # 保存为该页的 layout.json
                with open(os.path.join(page_dir, "layout.json"), "w", encoding="utf-8") as f:
                    json.dump(final_json_elements, f, indent=4, ensure_ascii=False)

            print(f"  [Done] 文件 {ppt_file} 处理完毕。")

        except Exception as e:
            print(f"  [Error] 处理文件 {ppt_file} 时出错: {e}")
            import traceback
            traceback.print_exc()

        finally:
            if os.path.exists(temp_worker_dir):
                shutil.rmtree(temp_worker_dir)
             # 如果需要批量处理所有文件，请注释掉这一行
            index+=1
            if index>100:
                break

if __name__ == "__main__":
    INPUT_DIR = "../pptOriginalData"
    OUTPUT_DIR = "../pptDataOutput"
    process_batch_ppts(INPUT_DIR, OUTPUT_DIR)